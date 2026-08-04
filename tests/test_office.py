"""Tests for the office (.xlsx / .pptx) tools, incl. undoability of binary edits."""

import pytest

pytest.importorskip("openpyxl")
pytest.importorskip("pptx")
# NOTE: python-docx is gated per-test (see the read_docx tests), not module-wide,
# so the xlsx/pptx suite still runs when only python-docx is missing.

from opendot.reversibility.engine import Reversibility
from opendot.reversibility.snapshots import IgnoreRules
from opendot.tools.local import Toolbox


@pytest.fixture(autouse=True)
def isolated_store(tmp_path, monkeypatch):
    monkeypatch.setenv("OPENDOT_HOME", str(tmp_path / "store"))
    yield


def _tb(tmp_path):
    wd = tmp_path / "ws"
    wd.mkdir()
    rev = Reversibility(workdir=str(wd), rules=IgnoreRules())
    return Toolbox(str(wd), reversibility=rev, confirm=lambda p: True), wd, rev


def _make_xlsx(path):
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "name"
    ws["B1"] = "score"
    ws["A2"] = "alice"
    ws["B2"] = 100
    wb.save(path)


def _make_pptx(path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.text = "Old Title"
    prs.save(path)


def _make_docx(path):
    from docx import Document

    doc = Document()
    doc.add_heading("Project Report", level=1)
    doc.add_paragraph("First paragraph of the body.")
    doc.add_paragraph("Second paragraph with details.")
    doc.save(path)


def test_office_tools_registered(tmp_path):
    tb, _, _ = _tb(tmp_path)
    names = {s["function"]["name"] for s in tb.specs()}
    assert {"read_xlsx", "edit_cell", "read_pptx", "edit_pptx_text"} <= names
    # read_docx is gated on python-docx being importable.
    import importlib.util

    if importlib.util.find_spec("docx") is not None:
        assert "read_docx" in names


def test_read_docx(tmp_path):
    pytest.importorskip("docx")
    tb, wd, _ = _tb(tmp_path)
    _make_docx(wd / "report.docx")

    out = tb.call("read_docx", {"path": "report.docx"})
    assert "Project Report" in out
    assert "First paragraph of the body." in out
    assert "Second paragraph with details." in out


def test_read_docx_missing_file(tmp_path):
    pytest.importorskip("docx")
    tb, _, _ = _tb(tmp_path)
    assert tb.call("read_docx", {"path": "nope.docx"}).startswith("error: file not found")


def test_read_and_edit_xlsx(tmp_path):
    tb, wd, _ = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")

    out = tb.call("read_xlsx", {"path": "data.xlsx"})
    assert "alice" in out and "score" in out

    res = tb.call("edit_cell", {"path": "data.xlsx", "cell": "B2", "value": "120"})
    assert "100" in res and "120" in res

    import openpyxl

    wb = openpyxl.load_workbook(wd / "data.xlsx")
    assert wb.active["B2"].value == 120  # coerced to int


def test_xlsx_missing_sheet_lists_available_sheets(tmp_path):
    tb, wd, rev = _tb(tmp_path)
    path = wd / "data.xlsx"
    _make_xlsx(path)

    read_result = tb.call("read_xlsx", {"path": "data.xlsx", "sheet": "Data"})
    edit_result = tb.call(
        "edit_cell",
        {"path": "data.xlsx", "cell": "B2", "value": "120", "sheet": "Data"},
    )

    expected = "error: no sheet 'Data'; sheets: Sheet"
    assert read_result == expected
    assert edit_result == expected
    assert rev.history() == []

    import openpyxl

    assert openpyxl.load_workbook(path).active["B2"].value == 100


def test_xlsx_edit_is_undoable(tmp_path):
    tb, wd, rev = _tb(tmp_path)
    _make_xlsx(wd / "data.xlsx")
    tb.call("edit_cell", {"path": "data.xlsx", "cell": "B2", "value": "999"})

    import openpyxl

    assert openpyxl.load_workbook(wd / "data.xlsx").active["B2"].value == 999
    rev.undo_last()  # restore the binary file exactly
    assert openpyxl.load_workbook(wd / "data.xlsx").active["B2"].value == 100


def test_read_and_edit_pptx(tmp_path):
    tb, wd, rev = _tb(tmp_path)
    _make_pptx(wd / "deck.pptx")

    out = tb.call("read_pptx", {"path": "deck.pptx"})
    assert "Old Title" in out

    res = tb.call(
        "edit_pptx_text",
        {"path": "deck.pptx", "find": "Old Title", "replace": "New Title"},
    )
    assert "1 run" in res

    from pptx import Presentation

    prs = Presentation(str(wd / "deck.pptx"))
    texts = [sh.text for sh in prs.slides[0].shapes if sh.has_text_frame]
    assert "New Title" in texts

    rev.undo_last()  # binary pptx restored
    prs2 = Presentation(str(wd / "deck.pptx"))
    texts2 = [sh.text for sh in prs2.slides[0].shapes if sh.has_text_frame]
    assert "Old Title" in texts2


def test_edit_pptx_missing_text_is_error(tmp_path):
    tb, wd, _ = _tb(tmp_path)
    _make_pptx(wd / "deck.pptx")
    res = tb.call("edit_pptx_text", {"path": "deck.pptx", "find": "nope", "replace": "x"})
    assert "not found" in res


def test_office_edit_outside_workspace_is_irreversible(tmp_path):
    """An office edit to a file outside the working dir isn't covered by the
    snapshot, so it's recorded as not-undoable (undo that doesn't lie)."""
    wd = tmp_path / "ws"
    wd.mkdir()
    rev = Reversibility(workdir=str(wd), rules=IgnoreRules())
    tb = Toolbox(str(wd), reversibility=rev, confirm=lambda p: True)

    outside = tmp_path / "outside.xlsx"  # sibling of ws/, not under it
    _make_xlsx(outside)

    res = tb.call("edit_cell", {"path": str(outside), "cell": "B2", "value": "999"})
    assert "999" in res
    last = rev.history()[-1]
    assert last.reversible is False
    assert "not undoable" in last.note


def test_office_edit_outside_declined_is_skipped(tmp_path):
    wd = tmp_path / "ws"
    wd.mkdir()
    rev = Reversibility(workdir=str(wd), rules=IgnoreRules())
    tb = Toolbox(str(wd), reversibility=rev, confirm=lambda p: False)  # decline

    outside = tmp_path / "outside.xlsx"
    _make_xlsx(outside)
    import openpyxl

    before = openpyxl.load_workbook(outside).active["B2"].value
    res = tb.call("edit_cell", {"path": str(outside), "cell": "B2", "value": "999"})
    assert res.startswith("skipped")
    after = openpyxl.load_workbook(outside).active["B2"].value
    assert after == before  # file on disk unchanged


def test_read_pptx_respects_max_slides(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    tb, wd, _ = _tb(tmp_path)
    prs = Presentation()
    for n in range(1, 6):  # 5 slides
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        shape.text_frame.text = f"Slide {n} body"
    prs.save(wd / "big.pptx")

    out = tb.call("read_pptx", {"path": "big.pptx", "max_slides": 3})
    assert "Slide 3 body" in out
    assert "Slide 4 body" not in out
    assert "... (2 more slides)" in out

    full = tb.call("read_pptx", {"path": "big.pptx"})  # default cap is 50
    assert "Slide 5 body" in full
    assert "more slides" not in full


def test_read_pptx_negative_cap_clamps_to_none(tmp_path):
    """A negative cap means 'show none', not garbage arithmetic (e.g. '6 more' for
    a 5-slide deck). It clamps to 0, so the count in the notice stays correct."""
    from pptx import Presentation
    from pptx.util import Inches

    tb, wd, _ = _tb(tmp_path)
    prs = Presentation()
    for n in range(1, 6):  # 5 slides
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        shape.text_frame.text = f"Slide {n} body"
    prs.save(wd / "big.pptx")

    out = tb.call("read_pptx", {"path": "big.pptx", "max_slides": -1})
    assert "Slide 1 body" not in out  # nothing shown
    assert "... (5 more slides)" in out  # correct count, not "6 more"


def test_read_docx_respects_max_paragraphs(tmp_path):
    pytest.importorskip("docx")
    from docx import Document

    tb, wd, _ = _tb(tmp_path)
    doc = Document()
    for n in range(1, 11):  # 10 paragraphs
        doc.add_paragraph(f"Paragraph {n}")
    doc.save(wd / "big.docx")

    out = tb.call("read_docx", {"path": "big.docx", "max_paragraphs": 4})
    assert "Paragraph 4" in out
    assert "Paragraph 5" not in out
    assert "... (6 more paragraphs)" in out

    full = tb.call("read_docx", {"path": "big.docx"})  # default cap is 200
    assert "Paragraph 10" in full
    assert "more paragraphs" not in full
